"""
PPT 单页生成主入口脚本
协调布局引擎、风格主题、图片管理器、图形渲染器、图表构建器、校验器。
完整流程：内容分析 → 布局选择 → 风格应用 → 图片任务创建 → 生图 → 延迟组装PPTX → 校验

适配：
- 星辰智能体 (TeleAgent) 平台：输出到 ./TeleAgent的工作空间/
- 跨平台兼容：Windows / macOS / Linux
- 断点续传：支持中断恢复
"""

import os
import sys
import json
import signal
import argparse
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 同级模块导入
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from layout_engine import (
    LayoutTemplate, Zone, ZoneType, TextAlign, VerticalAlign,
    get_layout, select_layout, SLIDE_WIDTH, SLIDE_HEIGHT,
)
from style_themes import (
    StyleTheme, ColorPalette, FontConfig,
    get_theme, get_auto_theme,
)
from image_manager import ImageManager, ImageTask, TaskStatus, ASPECT_RATIO_MAP
from graphics_renderer import (
    generate_gradient_background_fast, generate_overlay,
    generate_wave_decoration, generate_geometric_shapes,
    generate_number_circle, generate_horizontal_line,
    generate_vertical_line, generate_rounded_rect, generate_circle,
    hex_to_rgba,
)
from chart_builder import build_chart
from validator import validate_pptx


# ============================================================
# TeleAgent 工作空间适配
# ============================================================

# 优先使用环境变量，否则检测默认路径
TELEAGENT_WORKSPACE = os.environ.get("TELEAGENT_WORKSPACE", "")


def get_workspace_dir() -> str:
    """获取 TeleAgent 工作空间路径"""
    # 1. 环境变量
    if TELEAGENT_WORKSPACE and os.path.isdir(TELEAGENT_WORKSPACE):
        return TELEAGENT_WORKSPACE

    # 2. 默认路径（TeleAgent 平台约定）
    default_ws = os.path.join(".", "TeleAgent的工作空间")
    if os.path.isdir(default_ws):
        return default_ws

    # 3. 当前目录
    return "."


def resolve_output_path(filename: str) -> str:
    """解析输出文件路径"""
    ws = get_workspace_dir()
    path = os.path.join(ws, filename)
    os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
    return path


def resolve_image_dir() -> str:
    """解析图片输出目录"""
    ws = get_workspace_dir()
    img_dir = os.path.join(ws, "images")
    os.makedirs(img_dir, exist_ok=True)
    return img_dir


# ============================================================
# 中断恢复机制
# ============================================================

class StateManager:
    """进度状态管理，支持断点续传"""

    def __init__(self, state_file: str = ""):
        if not state_file:
            ws = get_workspace_dir()
            state_file = os.path.join(ws, ".ppt_gen_state.json")
        self.state_file = state_file
        self._setup_signal_handlers()

    def _setup_signal_handlers(self):
        """设置信号处理器，优雅退出时保存状态"""
        try:
            signal.signal(signal.SIGTERM, self._handle_exit)
            signal.signal(signal.SIGINT, self._handle_exit)
        except (OSError, ValueError):
            # 某些环境不支持信号处理（如线程中）
            pass

    def _handle_exit(self, signum, frame):
        print("\n⚠️ 收到中断信号，保存进度...")
        raise SystemExit(0)

    def save_state(self, data: dict):
        """保存当前进度"""
        try:
            with open(self.state_file, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def load_state(self) -> Optional[dict]:
        """加载上次中断的进度"""
        if os.path.exists(self.state_file):
            try:
                with open(self.state_file, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception:
                pass
        return None

    def clear_state(self):
        """清除已完成的状态文件"""
        if os.path.exists(self.state_file):
            try:
                os.remove(self.state_file)
            except Exception:
                pass


# ============================================================
# 内容分析
# ============================================================

@dataclass
class ContentAnalysis:
    """内容分析结果"""
    page_type: str                   # cover | content | data | mixed | comparison | timeline
    title: str = ""
    subtitle: str = ""
    section_number: str = ""
    items: list = field(default_factory=list)
    images: list = field(default_factory=list)
    charts: list = field(default_factory=list)
    stats: list = field(default_factory=list)
    has_chart: bool = False
    item_count: int = 0
    extra: dict = field(default_factory=dict)


def analyze_content(content: dict) -> ContentAnalysis:
    """分析输入内容结构"""
    page_type = content.get("page_type", "content")
    items = content.get("items", [])
    images = content.get("images", [])
    charts = content.get("charts", [])
    stats = content.get("stats", [])

    return ContentAnalysis(
        page_type=page_type,
        title=content.get("title", ""),
        subtitle=content.get("subtitle", ""),
        section_number=content.get("section_number", ""),
        items=items,
        images=images,
        charts=charts,
        stats=stats,
        has_chart=bool(charts) or bool(stats),
        item_count=len(items),
        extra={
            k: v for k, v in content.items()
            if k not in ("page_type", "title", "subtitle", "section_number",
                         "items", "images", "charts", "stats", "style", "layout")
        },
    )


# ============================================================
# 颜色辅助
# ============================================================

def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    return RGBColor(0x33, 0x33, 0x33)


def _get_color(palette: ColorPalette, key: str) -> RGBColor:
    color_hex = getattr(palette, key, "#333333")
    if color_hex.startswith("rgba"):
        color_hex = "#333333"
    return _hex_to_rgb(color_hex)


def _get_font_size(fonts: FontConfig, key: str) -> int:
    size_map = {
        "title": fonts.title_size_pt,
        "subtitle": fonts.subtitle_size_pt,
        "heading": fonts.heading_size_pt,
        "body": fonts.body_size_pt,
        "caption": fonts.caption_size_pt,
        "number": fonts.number_size_pt,
    }
    return size_map.get(key, fonts.body_size_pt)


# ============================================================
# 内容提取
# ============================================================

def _get_content_value(content: ContentAnalysis, content_key: str) -> Optional[str]:
    """根据 content_key 从内容分析结果中提取文本值"""
    if not content_key:
        return None

    # 直接属性
    if content_key in ("title", "subtitle", "section_number"):
        return getattr(content, content_key, None)

    # items[i].field 格式
    if content_key.startswith("items["):
        try:
            idx_str = content_key.split("[")[1].split("]")[0]
            idx = int(idx_str)
            rest = content_key.split("].", 1)
            if idx < len(content.items):
                item = content.items[idx]
                if len(rest) > 1:
                    field_name = rest[1]
                    return str(item.get(field_name, ""))
                else:
                    if isinstance(item, dict):
                        parts = []
                        if "number" in item:
                            parts.append(item["number"])
                        if "title" in item:
                            parts.append(item["title"])
                        if "description" in item:
                            parts.append(item["description"])
                        return "\n".join(parts) if parts else str(item)
                    return str(item)
        except (ValueError, IndexError):
            pass
        return None

    # stats[i].field 格式
    if content_key.startswith("stats["):
        try:
            idx_str = content_key.split("[")[1].split("]")[0]
            idx = int(idx_str)
            rest = content_key.split("].", 1)
            if idx < len(content.stats):
                stat = content.stats[idx]
                if len(rest) > 1:
                    return str(stat.get(rest[1], ""))
                return str(stat)
        except (ValueError, IndexError):
            pass
        return None

    # extra 字段
    if content_key in content.extra:
        return str(content.extra[content_key])

    return None


# ============================================================
# 幻灯片生成器
# ============================================================

class SlideGenerator:
    """PPT 单页生成器 — 延迟组装模式"""

    def __init__(
        self,
        content: dict,
        style: str = "auto",
        layout: str = "auto",
        output_path: str = "",
        image_dir: str = "",
    ):
        self.raw_content = content
        self.style_name = style
        self.layout_name = layout

        # TeleAgent 工作空间适配
        if not output_path:
            output_path = resolve_output_path("output.pptx")
        self.output_path = output_path

        if not image_dir:
            image_dir = resolve_image_dir()
        self.image_dir = image_dir

        self.analysis: Optional[ContentAnalysis] = None
        self.layout_template: Optional[LayoutTemplate] = None
        self.theme: Optional[StyleTheme] = None
        self.image_manager = ImageManager(image_dir)
        self.state_manager = StateManager()
        self.prs: Optional[Presentation] = None

    def generate(self) -> str:
        """执行完整的生成流程"""
        # 检查断点续传
        saved = self.state_manager.load_state()
        if saved and saved.get("phase", 0) >= 3:
            print("🔄 检测到中断状态，尝试恢复...")
            tasks_file = saved.get("tasks_file")
            if tasks_file and os.path.exists(tasks_file):
                self.image_manager.import_tasks(tasks_file)
                print(f"   恢复了 {len(self.image_manager.tasks)} 个图片任务")

        print("📋 Phase 1: 内容分析...")
        self.analysis = analyze_content(self.raw_content)
        print(f"   页面类型: {self.analysis.page_type}")
        print(f"   标题: {self.analysis.title}")
        print(f"   内容项数: {self.analysis.item_count}")

        print("\n📐 Phase 2: 布局选择...")
        if self.layout_name == "auto":
            self.layout_template = select_layout(
                self.analysis.page_type,
                self.analysis.item_count,
                self.analysis.has_chart,
            )
        else:
            self.layout_template = get_layout(self.layout_name)
        print(f"   布局: {self.layout_template.display_name}")

        print("\n🎨 Phase 3: 风格应用...")
        if self.style_name == "auto":
            self.theme = get_auto_theme(self.analysis.page_type)
        else:
            self.theme = get_theme(self.style_name)
        print(f"   风格: {self.theme.display_name}")

        print("\n🖼️  Phase 4: 创建图片任务...")
        self._create_image_tasks()

        # 保存状态（支持断点续传）
        tasks_file = self.image_manager.export_tasks()
        self.state_manager.save_state({
            "phase": 4,
            "tasks_file": tasks_file,
            "output_path": self.output_path,
        })

        print("\n🎨 Phase 5: 执行图片生成...")
        self._execute_image_tasks()

        print("\n📦 Phase 6: 组装 PPTX...")
        self._assemble()

        print("\n✅ Phase 7: 校验...")
        report = validate_pptx(self.output_path)
        report.print_report()

        # 清除状态文件
        self.state_manager.clear_state()

        print(f"\n🎉 输出文件: {self.output_path}")
        return self.output_path

    def _create_image_tasks(self):
        """为所有图片/图表/装饰区域创建任务（纯元数据，不生成文件）"""
        for zone in self.layout_template.zones:
            if zone.zone_type in (ZoneType.IMAGE, ZoneType.BACKGROUND,
                                   ZoneType.CHART, ZoneType.ICON):
                if zone.zone_type == ZoneType.CHART:
                    img_type = "chart"
                elif zone.zone_type == ZoneType.ICON:
                    img_type = "icon"
                else:
                    img_type = "photo"

                desc = ""
                if zone.content_key:
                    val = _get_content_value(self.analysis, zone.content_key)
                    if val:
                        desc = val

                for img_def in self.analysis.images:
                    if img_def.get("zone_id") == zone.id or img_def.get("type") == img_type:
                        desc = img_def.get("description", desc)
                        if img_def.get("aspect_ratio"):
                            zone.aspect_ratio = img_def["aspect_ratio"]
                        break

                self.image_manager.create_task(
                    zone_id=zone.id,
                    left=zone.left,
                    top=zone.top,
                    width=zone.width,
                    height=zone.height,
                    image_type=img_type,
                    description=desc,
                    style_prompt=self.theme.image_prompt_suffix,
                    preferred_ratio=zone.aspect_ratio,
                )

        summary = self.image_manager.get_summary()
        print(f"   创建了 {summary['total']} 个图片任务: {summary['by_type']}")

    def _execute_image_tasks(self):
        """执行所有图片生成任务"""
        # 1. 图标（直接 Pillow 绘制）
        for task in self.image_manager.get_by_type("icon"):
            if task.status != TaskStatus.PENDING:
                continue
            self._generate_icon(task)

        # 2. 图表（直接 Pillow 绘制）
        for task in self.image_manager.get_by_type("chart"):
            if task.status != TaskStatus.PENDING:
                continue
            self._generate_chart(task)

        # 3. 照片类 — 这里生成渐变占位，实际使用时通过生图工具替换
        for task in self.image_manager.get_by_type("photo"):
            if task.status != TaskStatus.PENDING:
                continue
            self._generate_photo_fallback(task)

        # 更新任务清单
        self.image_manager.export_tasks()
        summary = self.image_manager.get_summary()
        print(f"   图片生成完成: {summary['by_status']}")

    def _generate_icon(self, task: ImageTask):
        """生成图标"""
        self.image_manager.mark_generating(task.id)
        size = min(task.target_width, task.target_height)
        output_path = os.path.join(self.image_dir, f"{task.id}_icon.png")

        generate_number_circle(
            size=size,
            number="●",
            bg_color=self.theme.colors.primary,
            text_color=self.theme.colors.text_light,
            output_path=output_path,
        )
        self.image_manager.fill_image(task.id, output_path, crop_to_fit=False)

    def _generate_chart(self, task: ImageTask):
        """生成图表"""
        self.image_manager.mark_generating(task.id)
        output_path = os.path.join(self.image_dir, f"{task.id}_chart.png")

        # 检查是否为金字塔图元
        if "pyramid" in task.zone_id:
            from graphics_renderer import generate_pyramid_diagram
            generate_pyramid_diagram(
                task.target_width, task.target_height,
                levels=3,
                colors=[self.theme.colors.primary, self.theme.colors.secondary,
                        self.theme.colors.accent, "#CBDCE6"],
                output_path=output_path,
            )
            self.image_manager.fill_image(task.id, output_path, crop_to_fit=False)
            return

        # 检查是否为阶梯上升图元
        if "stairs" in task.zone_id:
            from graphics_renderer import generate_stepped_stairs
            generate_stepped_stairs(
                task.target_width, task.target_height,
                steps=4,
                colors=[self.theme.colors.primary, self.theme.colors.secondary,
                        self.theme.colors.accent, "#CBDCE6"],
                output_path=output_path,
            )
            self.image_manager.fill_image(task.id, output_path, crop_to_fit=False)
            return

        chart_config = None
        for chart in self.analysis.charts:
            chart_config = chart
            break
        for stat in self.analysis.stats:
            if stat.get("chart"):
                chart_config = stat["chart"]
                break

        if chart_config and isinstance(chart_config, dict):
            chart_type = chart_config.get("chart_type", "donut")
            build_chart(
                chart_type,
                task.target_width,
                task.target_height,
                output_path,
                value=chart_config.get("value", 68),
                color=self.theme.colors.primary,
                text_color=self.theme.colors.text_primary,
                label=chart_config.get("label", ""),
            )
        else:
            build_chart(
                "donut",
                task.target_width,
                task.target_height,
                output_path,
                value=68,
                color=self.theme.colors.primary,
                text_color=self.theme.colors.text_primary,
            )

        self.image_manager.fill_image(task.id, output_path, crop_to_fit=False)

    def _generate_photo_fallback(self, task: ImageTask):
        """
        照片类图片的备选渲染。
        在 TeleAgent 环境中，实际图片应由 AI 生图工具生成后调用
        image_manager.fill_image() 填充。此处生成渐变色备选图。
        """
        self.image_manager.mark_generating(task.id)
        output_path = os.path.join(self.image_dir, f"{task.id}_fallback.png")

        generate_gradient_background_fast(
            task.target_width,
            task.target_height,
            self.theme.colors.gradient_start,
            self.theme.colors.gradient_end,
            direction="diagonal",
            opacity=0.8,
            output_path=output_path,
        )
        self.image_manager.fill_image(task.id, output_path, crop_to_fit=False)

    def _assemble(self):
        """组装 PPTX — 延迟组装，只使用已完成的图片"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_WIDTH)
        self.prs.slide_height = Inches(SLIDE_HEIGHT)

        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(self.theme.colors.background)

        # 按 z_order 排序
        sorted_zones = sorted(self.layout_template.zones, key=lambda z: z.z_order)

        for zone in sorted_zones:
            if zone.zone_type == ZoneType.DECORATOR:
                self._add_decorator_to_slide(slide, zone)
            elif zone.zone_type in (ZoneType.IMAGE, ZoneType.BACKGROUND,
                                     ZoneType.CHART, ZoneType.ICON):
                self._add_image_to_slide(slide, zone)
            elif zone.zone_type in (ZoneType.TEXT, ZoneType.NUMBER):
                self._add_text_to_slide(slide, zone)

        # 确保输出目录存在
        out_dir = os.path.dirname(self.output_path)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)
        self.prs.save(self.output_path)
        print(f"   已保存: {self.output_path}")

    def _add_text_to_slide(self, slide, zone: Zone):
        text_content = _get_content_value(self.analysis, zone.content_key)
        if not text_content:
            return

        txBox = slide.shapes.add_textbox(
            Inches(zone.left), Inches(zone.top),
            Inches(zone.width), Inches(zone.height),
        )
        txBox.name = zone.id
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = text_content.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line

            if zone.text_align == TextAlign.CENTER:
                p.alignment = PP_ALIGN.CENTER
            elif zone.text_align == TextAlign.RIGHT:
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            font_size = _get_font_size(self.theme.fonts, zone.font_size_key)
            color = _get_color(self.theme.colors, zone.color_key)

            run = p.runs[0] if p.runs else p.add_run()
            if not p.runs:
                run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = self.theme.fonts.title_font

            if zone.is_bold is not None:
                run.font.bold = zone.is_bold
            elif zone.font_size_key in ("title", "heading", "number"):
                run.font.bold = self.theme.fonts.title_bold

    def _add_image_to_slide(self, slide, zone: Zone):
        """添加图片 — 只使用已完成(done)的任务"""
        task_id = f"img_{zone.id}"
        task = self.image_manager.get_task(task_id)

        if task and task.status == TaskStatus.DONE and task.image_path:
            if os.path.exists(task.image_path):
                try:
                    slide.shapes.add_picture(
                        task.image_path,
                        Inches(zone.left), Inches(zone.top),
                        Inches(zone.width), Inches(zone.height),
                    )
                except Exception as e:
                    print(f"   ⚠️ 添加图片失败 ({zone.id}): {e}")

    def _add_decorator_to_slide(self, slide, zone: Zone):
        """添加装饰元素"""
        dpi = 96
        w = max(1, int(zone.width * dpi))
        h = max(1, int(zone.height * dpi))
        output_path = os.path.join(self.image_dir, f"dec_{zone.id}.png")

        dec = self.theme.decorators
        from graphics_renderer import generate_radial_hub

        if zone.id == "overlay":
            generate_overlay(w, h, self.theme.colors.background, 0.6,
                            output_path=output_path)
        elif zone.id == "center_circle":
            generate_circle(min(w, h),
                           border_color=self.theme.colors.primary,
                           border_width=3,
                           output_path=output_path)
        elif zone.id in ("center_hub", "orbit_visual"):
            generate_radial_hub(w, h, node_count=4, center_label="?",
                                primary_color=self.theme.colors.primary,
                                output_path=output_path)
        elif zone.id == "center_box":
            generate_rounded_rect(w, h, bg_color=self.theme.colors.surface,
                                  border_color=self.theme.colors.border,
                                  border_width=2, radius=8,
                                  output_path=output_path)
        elif zone.id in ("part_badge", "left_bar", "goal_circle", "panel_bg"):
            generate_rounded_rect(w, h, bg_color=self.theme.colors.primary,
                                  radius=6, output_path=output_path)
        elif zone.id == "timeline_line":
            generate_horizontal_line(w, h, self.theme.colors.primary, 3,
                                    output_path=output_path)
        elif zone.id == "divider":
            generate_vertical_line(w, h, self.theme.colors.border, 2,
                                  output_path=output_path)
        elif "dot" in zone.id:
            generate_number_circle(min(w, h), "•",
                                  self.theme.colors.primary,
                                  self.theme.colors.text_light,
                                  output_path=output_path)
        elif dec.has_wave_decoration:
            generate_wave_decoration(w, h, self.theme.colors.primary, 0.15,
                                    output_path=output_path)
        elif dec.has_geometric_shapes:
            generate_geometric_shapes(w, h, self.theme.colors.primary, 0.08,
                                     output_path=output_path)
        else:
            generate_gradient_background_fast(
                w, h, self.theme.colors.gradient_start,
                self.theme.colors.gradient_end, opacity=0.3,
                output_path=output_path)

        if os.path.exists(output_path):
            try:
                slide.shapes.add_picture(
                    output_path,
                    Inches(zone.left), Inches(zone.top),
                    Inches(zone.width), Inches(zone.height),
                )
            except Exception as e:
                print(f"   ⚠️ 装饰元素添加失败 ({zone.id}): {e}")


# ============================================================
# CLI 入口
# ============================================================

def main():
    parser = argparse.ArgumentParser(description="PPT 单页智能生成")
    parser.add_argument("--input", "-i", required=True, help="输入 JSON 文件路径")
    parser.add_argument("--output", "-o", default="", help="输出 PPTX 文件路径（默认: TeleAgent工作空间）")
    parser.add_argument("--style", "-s", default="auto",
                        help="风格名称 (auto|business_blue|nature_green|tech_dark|...)")
    parser.add_argument("--layout", "-l", default="auto",
                        help="布局名称 (auto|cover_fullimage|content_grid_2x3|...)")
    parser.add_argument("--image-dir", default="", help="图片输出目录（默认: TeleAgent工作空间/images）")
    parser.add_argument("--resume", action="store_true", help="从上次中断处恢复")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        content = json.load(f)

    generator = SlideGenerator(
        content=content,
        style=args.style if args.style != "auto" else content.get("style", "auto"),
        layout=args.layout if args.layout != "auto" else content.get("layout", "auto"),
        output_path=args.output if args.output else "",
        image_dir=args.image_dir if args.image_dir else "",
    )
    generator.generate()


if __name__ == "__main__":
    main()
