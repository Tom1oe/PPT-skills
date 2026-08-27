"""
校验器模块
对生成的 PPTX 进行质量检查：文字重叠、图片错位、比例形变。
"""

import os
import sys
import json
import argparse
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image


# 幻灯片标准尺寸（EMU单位）
SLIDE_WIDTH_EMU = 12192000   # 13.333 inches
SLIDE_HEIGHT_EMU = 6858000   # 7.5 inches

# 阈值配置
TEXT_OVERLAP_IOU_THRESHOLD = 0.05   # 文本框重叠 IoU 阈值
POSITION_DEVIATION_INCHES = 0.1     # 位置偏差阈值（英寸）
ASPECT_RATIO_TOLERANCE = 0.02       # 比例偏差容忍度 (2%)
# 估算每行字符数的基准（12pt 字体约 0.1 英寸/字符）
CHAR_WIDTH_INCHES = 0.1
LINE_HEIGHT_FACTOR = 1.4


@dataclass
class BoundingBox:
    """边界框"""
    left: float    # 英寸
    top: float     # 英寸
    width: float   # 英寸
    height: float  # 英寸
    label: str = ""

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def area(self) -> float:
        return self.width * self.height


@dataclass
class ValidationIssue:
    """校验问题"""
    level: str       # error | warning | suggestion
    category: str    # text_overlap | image_misalign | aspect_distortion
    message: str
    element_id: str = ""
    details: dict = field(default_factory=dict)


@dataclass
class ValidationReport:
    """校验报告"""
    passed: bool = True
    issues: list[ValidationIssue] = field(default_factory=list)

    @property
    def errors(self) -> list[ValidationIssue]:
        return [i for i in self.issues if i.level == "error"]

    @property
    def warnings(self) -> list[ValidationIssue]:
        return [i for i in self.issues if i.level == "warning"]

    @property
    def suggestions(self) -> list[ValidationIssue]:
        return [i for i in self.issues if i.level == "suggestion"]

    def add(self, issue: ValidationIssue):
        self.issues.append(issue)
        if issue.level == "error":
            self.passed = False

    def to_dict(self) -> dict:
        return {
            "passed": self.passed,
            "total_issues": len(self.issues),
            "errors": [
                {"category": i.category, "message": i.message,
                 "element": i.element_id, "details": i.details}
                for i in self.errors
            ],
            "warnings": [
                {"category": i.category, "message": i.message,
                 "element": i.element_id, "details": i.details}
                for i in self.warnings
            ],
            "suggestions": [
                {"category": i.category, "message": i.message,
                 "element": i.element_id, "details": i.details}
                for i in self.suggestions
            ],
        }

    def print_report(self):
        status = "✅ 通过" if self.passed else "❌ 未通过"
        print(f"\n{'='*60}")
        print(f"  校验结果: {status}")
        print(f"  总问题数: {len(self.issues)}")
        print(f"  错误: {len(self.errors)}  警告: {len(self.warnings)}  建议: {len(self.suggestions)}")
        print(f"{'='*60}")
        for i in self.errors:
            print(f"  ❌ [{i.category}] {i.message}")
        for i in self.warnings:
            print(f"  ⚠️ [{i.category}] {i.message}")
        for i in self.suggestions:
            print(f"  💡 [{i.category}] {i.message}")
        print()


def _emu_to_inches(emu: int) -> float:
    """EMU 转英寸"""
    return emu / 914400


def _get_shape_bbox(shape) -> BoundingBox:
    """获取 shape 的边界框（英寸）"""
    return BoundingBox(
        left=_emu_to_inches(shape.left),
        top=_emu_to_inches(shape.top),
        width=_emu_to_inches(shape.width),
        height=_emu_to_inches(shape.height),
        label=shape.name or "",
    )


def _calc_iou(a: BoundingBox, b: BoundingBox) -> float:
    """计算两个边界框的 IoU（交并比）"""
    inter_left = max(a.left, b.left)
    inter_top = max(a.top, b.top)
    inter_right = min(a.right, b.right)
    inter_bottom = min(a.bottom, b.bottom)

    if inter_right <= inter_left or inter_bottom <= inter_top:
        return 0.0

    inter_area = (inter_right - inter_left) * (inter_bottom - inter_top)
    union_area = a.area + b.area - inter_area
    if union_area <= 0:
        return 0.0
    return inter_area / union_area


def _estimate_text_overflow(shape) -> Optional[dict]:
    """估算文本是否溢出文本框"""
    if not shape.has_text_frame:
        return None

    bbox = _get_shape_bbox(shape)
    text = shape.text_frame.text
    if not text.strip():
        return None

    # 估算所需面积
    total_chars = len(text)
    # 获取字体大小
    font_size_pt = 12  # 默认
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                font_size_pt = run.font.size.pt
                break

    char_w = font_size_pt / 72 * 0.6  # 大致字符宽度（英寸）
    line_h = font_size_pt / 72 * LINE_HEIGHT_FACTOR
    chars_per_line = max(1, int(bbox.width / char_w))
    needed_lines = max(1, (total_chars + chars_per_line - 1) // chars_per_line)
    needed_height = needed_lines * line_h

    if needed_height > bbox.height * 1.1:
        return {
            "chars": total_chars,
            "needed_height": round(needed_height, 2),
            "available_height": round(bbox.height, 2),
            "overflow_ratio": round(needed_height / bbox.height, 2),
        }
    return None


# ============================================================
# 校验函数
# ============================================================

def check_text_overlap(slide, report: ValidationReport):
    """检测文本框之间的重叠"""
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            box_a = _get_shape_bbox(text_shapes[i])
            box_b = _get_shape_bbox(text_shapes[j])
            iou = _calc_iou(box_a, box_b)
            if iou > TEXT_OVERLAP_IOU_THRESHOLD:
                report.add(ValidationIssue(
                    level="error",
                    category="text_overlap",
                    message=f"文本框 '{box_a.label}' 与 '{box_b.label}' 重叠 (IoU={iou:.2%})",
                    element_id=f"{box_a.label}+{box_b.label}",
                    details={"iou": round(iou, 4)},
                ))
            elif iou > 0:
                report.add(ValidationIssue(
                    level="warning",
                    category="text_overlap",
                    message=f"文本框 '{box_a.label}' 与 '{box_b.label}' 轻微重叠 (IoU={iou:.2%})",
                    element_id=f"{box_a.label}+{box_b.label}",
                    details={"iou": round(iou, 4)},
                ))

    # 检查文本溢出
    for s in text_shapes:
        overflow = _estimate_text_overflow(s)
        if overflow:
            bbox = _get_shape_bbox(s)
            if overflow["overflow_ratio"] > 1.5:
                report.add(ValidationIssue(
                    level="error",
                    category="text_overflow",
                    message=f"文本框 '{bbox.label}' 内容严重溢出 "
                            f"(需要 {overflow['needed_height']}\", 可用 {overflow['available_height']}\")",
                    element_id=bbox.label,
                    details=overflow,
                ))
            else:
                report.add(ValidationIssue(
                    level="warning",
                    category="text_overflow",
                    message=f"文本框 '{bbox.label}' 内容可能溢出",
                    element_id=bbox.label,
                    details=overflow,
                ))


def check_image_alignment(slide, report: ValidationReport):
    """检测图片是否超出幻灯片范围或错位"""
    slide_w = _emu_to_inches(SLIDE_WIDTH_EMU)
    slide_h = _emu_to_inches(SLIDE_HEIGHT_EMU)

    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            bbox = _get_shape_bbox(shape)

            # 检查是否超出边界
            if bbox.left < -POSITION_DEVIATION_INCHES:
                report.add(ValidationIssue(
                    level="error",
                    category="image_misalign",
                    message=f"图片 '{bbox.label}' 左侧超出幻灯片边界 ({bbox.left:.2f}\")",
                    element_id=bbox.label,
                ))
            if bbox.top < -POSITION_DEVIATION_INCHES:
                report.add(ValidationIssue(
                    level="error",
                    category="image_misalign",
                    message=f"图片 '{bbox.label}' 上方超出幻灯片边界 ({bbox.top:.2f}\")",
                    element_id=bbox.label,
                ))
            if bbox.right > slide_w + POSITION_DEVIATION_INCHES:
                report.add(ValidationIssue(
                    level="error",
                    category="image_misalign",
                    message=f"图片 '{bbox.label}' 右侧超出幻灯片边界 "
                            f"(right={bbox.right:.2f}\", slide_w={slide_w:.2f}\")",
                    element_id=bbox.label,
                ))
            if bbox.bottom > slide_h + POSITION_DEVIATION_INCHES:
                report.add(ValidationIssue(
                    level="warning",
                    category="image_misalign",
                    message=f"图片 '{bbox.label}' 下方超出幻灯片边界 "
                            f"(bottom={bbox.bottom:.2f}\", slide_h={slide_h:.2f}\")",
                    element_id=bbox.label,
                ))

            # 检查尺寸是否过小
            if bbox.width < 0.3 or bbox.height < 0.3:
                report.add(ValidationIssue(
                    level="warning",
                    category="image_misalign",
                    message=f"图片 '{bbox.label}' 尺寸过小 ({bbox.width:.2f}\"×{bbox.height:.2f}\")",
                    element_id=bbox.label,
                ))


def check_aspect_ratio(slide, image_dir: Optional[str] = None,
                       report: Optional[ValidationReport] = None):
    """检测图片嵌入比例是否有形变"""
    if report is None:
        report = ValidationReport()

    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            bbox = _get_shape_bbox(shape)
            embed_ratio = bbox.width / bbox.height if bbox.height > 0 else 1.0

            # 尝试获取原始图片比例
            try:
                image = shape.image
                blob = image.blob
                from io import BytesIO
                img = Image.open(BytesIO(blob))
                orig_w, orig_h = img.size
                orig_ratio = orig_w / orig_h if orig_h > 0 else 1.0

                deviation = abs(embed_ratio - orig_ratio) / orig_ratio
                if deviation > ASPECT_RATIO_TOLERANCE:
                    report.add(ValidationIssue(
                        level="error" if deviation > 0.1 else "warning",
                        category="aspect_distortion",
                        message=f"图片 '{bbox.label}' 存在比例形变 "
                                f"(原始比例 {orig_ratio:.3f}, 嵌入比例 {embed_ratio:.3f}, "
                                f"偏差 {deviation:.1%})",
                        element_id=bbox.label,
                        details={
                            "original_ratio": round(orig_ratio, 3),
                            "embed_ratio": round(embed_ratio, 3),
                            "deviation": round(deviation, 4),
                        },
                    ))
            except Exception:
                # 无法读取原始图片时跳过
                report.add(ValidationIssue(
                    level="suggestion",
                    category="aspect_distortion",
                    message=f"无法检测图片 '{bbox.label}' 的原始比例",
                    element_id=bbox.label,
                ))

    return report


# ============================================================
# 主入口
# ============================================================

def validate_pptx(pptx_path: str, slide_index: int = 0) -> ValidationReport:
    """
    验证 PPTX 文件的第 slide_index 页。
    返回 ValidationReport。
    """
    if not os.path.exists(pptx_path):
        report = ValidationReport(passed=False)
        report.add(ValidationIssue(
            level="error", category="file",
            message=f"文件不存在: {pptx_path}",
        ))
        return report

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        report = ValidationReport(passed=False)
        report.add(ValidationIssue(
            level="error", category="file",
            message=f"幻灯片索引 {slide_index} 超出范围 (共 {len(prs.slides)} 页)",
        ))
        return report

    slide = prs.slides[slide_index]
    report = ValidationReport()

    # 1. 文字重叠检测
    check_text_overlap(slide, report)

    # 2. 图片错位检测
    check_image_alignment(slide, report)

    # 3. 比例形变检测
    check_aspect_ratio(slide, report=report)

    return report


def main():
    parser = argparse.ArgumentParser(description="校验 PPTX 文件质量")
    parser.add_argument("--input", "-i", required=True, help="PPTX 文件路径")
    parser.add_argument("--slide", "-s", type=int, default=0, help="幻灯片索引（默认0）")
    parser.add_argument("--json", action="store_true", help="以 JSON 格式输出结果")
    args = parser.parse_args()

    report = validate_pptx(args.input, args.slide)

    if args.json:
        print(json.dumps(report.to_dict(), ensure_ascii=False, indent=2))
    else:
        report.print_report()

    sys.exit(0 if report.passed else 1)


if __name__ == "__main__":
    main()
